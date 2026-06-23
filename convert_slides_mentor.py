import os
import re
import sys
import subprocess

# Ensure python-pptx is installed
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError:
    print("python-pptx not found. Installing python-pptx...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

# Color Palette (Dark Tech Theme)
BG_COLOR = RGBColor(12, 16, 31)         # Deep slate/navy
TEXT_WHITE = RGBColor(255, 255, 255)    # White
TEXT_LIGHT_GRAY = RGBColor(216, 222, 233) # Content gray
ACCENT_BLUE = RGBColor(0, 210, 255)     # Neon electric blue
ACCENT_PURPLE = RGBColor(190, 80, 255)   # Neon purple

def parse_markdown_slides(filepath):
    with open(filepath, 'r', encoding='utf-8') as f:
        content = f.read()

    # Split by horizontal rule separator (---)
    slides_raw = re.split(r'\n---\n', content)
    slides_data = []

    for slide_text in slides_raw:
        slide_text = slide_text.strip()
        if not slide_text:
            continue

        # Skip the top document file title/header that is not a slide
        if not re.search(r'## Slide \d+', slide_text, re.IGNORECASE):
            continue

        lines = slide_text.split('\n')
        slide_info = {
            'header': '',
            'title': '',
            'subtitle': '',
            'visual_suggestions': [],
            'presenter_script': [],
            'content_bullets': [] # List of {'text': str, 'level': int}
        }

        current_section = None

        for line in lines:
            line_stripped = line.strip()
            if not line_stripped:
                continue

            # Header slide match
            if line_stripped.startswith('## '):
                slide_info['header'] = line_stripped.replace('## ', '').strip()
                current_section = None
                continue
            elif line_stripped.startswith('# '):
                slide_info['header'] = line_stripped.replace('# ', '').strip()
                current_section = None
                continue

            # Keys matching
            slide_title_match = re.match(r'^\*\s+\*\*Slide Title\*\*:\s*(.*)', line_stripped, re.IGNORECASE)
            if slide_title_match:
                slide_info['title'] = slide_title_match.group(1).strip()
                current_section = None
                continue

            subtitle_match = re.match(r'^\*\s+\*\*Subtitle\*\*:\s*(.*)', line_stripped, re.IGNORECASE)
            if subtitle_match:
                slide_info['subtitle'] = subtitle_match.group(1).strip()
                current_section = None
                continue

            visual_match = re.match(r'^\*\s+\*\*Visual Suggestion\*\*:\s*(.*)', line_stripped, re.IGNORECASE)
            if visual_match:
                val = visual_match.group(1).strip()
                if val:
                    slide_info['visual_suggestions'].append(val)
                current_section = 'visual'
                continue

            script_match = re.match(r'^\*\s+\*\*Presenter Script\*\*:\s*(.*)', line_stripped, re.IGNORECASE)
            if script_match:
                val = script_match.group(1).strip()
                if val:
                    val = re.sub(r'^>\s*', '', val)
                    slide_info['presenter_script'].append(val)
                current_section = 'script'
                continue

            # Reset section if a new primary bullet key is encountered
            other_bullet_match = re.match(r'^\*\s+\*\*(.*?)\*\*:\s*(.*)', line_stripped)
            if other_bullet_match:
                current_section = None

            # Check if this line is part of visual suggestions or presenter script
            if current_section == 'visual':
                cleaned = re.sub(r'^[\*\-\+]\s*', '', line_stripped)
                slide_info['visual_suggestions'].append(cleaned)
                continue
            elif current_section == 'script':
                cleaned = re.sub(r'^>\s*', '', line_stripped)
                cleaned = re.sub(r'^[\*\-\+]\s*', '', cleaned)
                slide_info['presenter_script'].append(cleaned)
                continue

            # Otherwise, it's a slide content bullet point
            # Determine bullet level based on indentation (leading spaces)
            leading_spaces = len(line) - len(line.lstrip())
            level = 0 if leading_spaces < 2 else 1
            
            # Clean up the bullet marker
            cleaned = re.sub(r'^[\*\-\+]\s*', '', line_stripped)
            # Remove markdown checklist markers [ ] if present
            cleaned = re.sub(r'^\[\s*\]\s*', '', cleaned)
            
            slide_info['content_bullets'].append({
                'text': cleaned,
                'level': level
            })

        slide_info['title'] = slide_info['title'] or slide_info['header']
        slides_data.append(slide_info)

    return slides_data

def add_markdown_runs(paragraph, text, default_color, default_size, default_font="Segoe UI"):
    # Check if text contains special headers to apply accent colors
    prefix = ""
    rest = text
    
    if text.startswith("Interactive Question:"):
        prefix = "Interactive Question:"
        rest = text[len("Interactive Question:"):]
        run = paragraph.add_run()
        run.text = prefix
        run.font.name = default_font
        run.font.size = default_size
        run.font.bold = True
        run.font.color.rgb = ACCENT_PURPLE
    elif text.startswith("Interactive Checkpoint:"):
        prefix = "Interactive Checkpoint:"
        rest = text[len("Interactive Checkpoint:"):]
        run = paragraph.add_run()
        run.text = prefix
        run.font.name = default_font
        run.font.size = default_size
        run.font.bold = True
        run.font.color.rgb = ACCENT_PURPLE
    elif text.startswith("Traditional Coding (Toaster):"):
        prefix = "Traditional Coding (Toaster):"
        rest = text[len("Traditional Coding (Toaster):"):]
        run = paragraph.add_run()
        run.text = prefix
        run.font.name = default_font
        run.font.size = default_size
        run.font.bold = True
        run.font.color.rgb = ACCENT_BLUE
    elif text.startswith("AI Coding (Smart Robot):"):
        prefix = "AI Coding (Smart Robot):"
        rest = text[len("AI Coding (Smart Robot):"):]
        run = paragraph.add_run()
        run.text = prefix
        run.font.name = default_font
        run.font.size = default_size
        run.font.bold = True
        run.font.color.rgb = ACCENT_BLUE

    # Split the remaining text by '**' for bold formatting
    parts = rest.split('**')
    for idx, part in enumerate(parts):
        if not part:
            continue
        run = paragraph.add_run()
        run.text = part
        run.font.name = default_font
        run.font.size = default_size
        run.font.bold = (idx % 2 == 1)
        
        # Color matching
        if idx % 2 == 1:
            run.font.color.rgb = TEXT_WHITE
        else:
            run.font.color.rgb = default_color

def apply_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def create_presentation(slides_data, assets_dir, output_path):
    prs = Presentation()
    
    # 16:9 Widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Use blank layout 6
    blank_layout = prs.slide_layouts[6]

    for idx, slide_info in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        apply_background(slide)
        
        slide_num = idx + 1
        
        # --- Left Column: Text & Content ---
        if slide_num == 1:
            # Decorative Purple Line
            line = slide.shapes.add_shape(
                1, # MSO_SHAPE.RECTANGLE
                Inches(0.8), Inches(1.8), Inches(0.12), Inches(3.6)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = ACCENT_PURPLE
            line.line.color.rgb = ACCENT_PURPLE
            
            # Text box
            txBox = slide.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(6.0), Inches(3.6))
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0)
            
            # Title
            p_title = tf.paragraphs[0]
            p_title.text = slide_info['title']
            p_title.font.name = "Segoe UI"
            p_title.font.size = Pt(44)
            p_title.font.bold = True
            p_title.font.color.rgb = TEXT_WHITE
            p_title.space_after = Pt(20)
            
            # Subtitle
            if slide_info['subtitle']:
                p_sub = tf.add_paragraph()
                p_sub.text = slide_info['subtitle']
                p_sub.font.name = "Segoe UI"
                p_sub.font.size = Pt(20)
                p_sub.font.color.rgb = ACCENT_BLUE
                
        else:
            # Normal slide
            # Header
            txHeader = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8))
            tfHeader = txHeader.text_frame
            tfHeader.word_wrap = True
            tfHeader.margin_left = tfHeader.margin_top = tfHeader.margin_right = tfHeader.margin_bottom = Inches(0)
            
            p_head = tfHeader.paragraphs[0]
            p_head.text = slide_info['title']
            p_head.font.name = "Segoe UI"
            p_head.font.size = Pt(28)
            p_head.font.bold = True
            p_head.font.color.rgb = TEXT_WHITE
            
            # Horizontal line separator
            line = slide.shapes.add_shape(
                1, # RECTANGLE
                Inches(0.8), Inches(1.2), Inches(11.733), Inches(0.02)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = ACCENT_BLUE
            line.line.color.rgb = ACCENT_BLUE
            
            # Main Left Body content (expanded to 6.3 inches width for teaching points)
            txBody = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(6.3), Inches(5.0))
            tf = txBody.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0)
            
            first = True
            for item in slide_info['content_bullets']:
                if first:
                    p = tf.paragraphs[0]
                    first = False
                else:
                    p = tf.add_paragraph()
                
                p.level = item['level']
                p.space_after = Pt(8)
                
                # Apply custom fonts, colors, and formatting based on level
                if p.level == 0:
                    bullet_prefix = ""
                    # Check if it starts with interactive checkpoint
                    if "Interactive Question:" in item['text'] or "Interactive Checkpoint:" in item['text']:
                        bullet_prefix = ""
                    else:
                        bullet_prefix = "• "
                    
                    add_markdown_runs(p, f"{bullet_prefix}{item['text']}", TEXT_WHITE, Pt(16))
                else:
                    add_markdown_runs(p, f"  - {item['text']}", TEXT_LIGHT_GRAY, Pt(14))

        # --- Right Column: Image block ---
        image_path = os.path.join(assets_dir, f"slide_{slide_num}.png")
        if os.path.exists(image_path):
            left = Inches(7.5)
            top = Inches(1.6)
            width = Inches(5.0)
            height = Inches(4.5)
            slide.shapes.add_picture(image_path, left, top, width, height)

        # --- Slide Notes (Presenter Script Backup) ---
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        
        notes_content = []
        if slide_info['visual_suggestions']:
            notes_content.append("--- VISUAL SUGGESTIONS ---")
            notes_content.extend([f"- {s}" for s in slide_info['visual_suggestions']])
            notes_content.append("")
            
        if slide_info['presenter_script']:
            notes_content.append("--- PRESENTER SCRIPT ---")
            notes_content.extend(slide_info['presenter_script'])
            notes_content.append("")

        notes_text_frame.text = "\n".join(notes_content).strip()

    try:
        prs.save(output_path)
        print(f"Successfully generated PowerPoint presentation at: {output_path}")
    except PermissionError:
        print("\n" + "="*60)
        print(f"ERROR: Permission denied when saving to '{output_path}'.")
        print("This usually means the file is open in Microsoft PowerPoint or another application.")
        print("Please CLOSE the PowerPoint window and run this script again!")
        print("="*60 + "\n")
        sys.exit(1)

if __name__ == "__main__":
    input_file = os.path.join("module_01_intro_to_ai", "slides_mentor.md")
    output_file = os.path.join("module_01_intro_to_ai", "slides_mentor.pptx")
    assets_folder = os.path.join("module_01_intro_to_ai", "assets")
    
    if not os.path.exists(input_file):
        print(f"Error: Could not find input file: {input_file}")
        sys.exit(1)
        
    slides_data = parse_markdown_slides(input_file)
    create_presentation(slides_data, assets_folder, output_file)
