import os
import re
import sys
import shutil
import subprocess

# Ensure python-pptx is installed
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("python-pptx not found. Installing python-pptx...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

# Color Palette (Dark Tech Theme)
BG_COLOR = RGBColor(12, 16, 31)         # Deep slate/navy
TEXT_WHITE = RGBColor(255, 255, 255)    # White
TEXT_LIGHT_GRAY = RGBColor(216, 222, 233) # Content gray
ACCENT_BLUE = RGBColor(0, 210, 255)     # Neon electric blue
ACCENT_PURPLE = RGBColor(190, 80, 255)   # Neon purple

# Image mapping from the artifact directory
ARTIFACT_DIR = r"C:\Users\User\.gemini\antigravity-ide\brain\e47cdacf-21d9-4fc8-b2d6-d8e353f2d9fc"
IMAGE_MAPPING = {
    1: "slide1_robot_1782041302224.png",
    2: "slide2_chore_robot_1782041314282.png",
    3: "slide3_toaster_robot_1782041326343.png",
    4: "slide4_recipe_1782041339928.png",
    5: "slide5_chef_1782041356358.png",
    6: "slide6_sorting_1782041369320.png",
    7: "slide7_flowchart_1782041381183.png",
    8: "slide8_teachable_1782041393942.png",
    9: "slide9_daily_ai_1782041408378.png",
    10: "slide10_challenge_1782041422572.png",
    11: "slide11_alive_1782044904562.png",
    12: "slide12_takeover_1782044916304.png",
    13: "slide13_chatgpt_1782044930890.png",
    14: "slide14_mistake_1782044944705.png",
    15: "slide15_genius_1782044960339.png"
}

def parse_markdown_slides(filepath):
    with open(filepath, 'r', encoding='utf-8') as f:
        content = f.read()

    # Split by horizontal rule separator (---)
    slides_raw = re.split(r'\n---\n', content)
    slides_data = []

    for slide_text in slides_raw:
        slide_text = slide_text.strip()
        if not slide_text:
            continue

        # Skip the top document file title/header that is not a slide
        if not re.search(r'## Slide \d+', slide_text, re.IGNORECASE):
            continue

        lines = slide_text.split('\n')
        slide_info = {
            'header': '',
            'title': '',
            'subtitle': '',
            'visual_suggestions': [],
            'presenter_script': [],
            'interactive_checkpoints': [],
            'other_bullets': []
        }

        current_section = None

        for line in lines:
            line_stripped = line.strip()
            if not line_stripped:
                continue

            # Header slide match
            if line_stripped.startswith('## '):
                slide_info['header'] = line_stripped.replace('## ', '').strip()
                current_section = None
                continue
            elif line_stripped.startswith('# '):
                slide_info['header'] = line_stripped.replace('# ', '').strip()
                current_section = None
                continue

            # Keys matching
            slide_title_match = re.match(r'^\*\s+\*\*Slide Title\*\*:\s*(.*)', line_stripped, re.IGNORECASE)
            if slide_title_match:
                slide_info['title'] = slide_title_match.group(1).strip()
                current_section = None
                continue

            subtitle_match = re.match(r'^\*\s+\*\*Subtitle\*\*:\s*(.*)', line_stripped, re.IGNORECASE)
            if subtitle_match:
                slide_info['subtitle'] = subtitle_match.group(1).strip()
                current_section = None
                continue

            visual_match = re.match(r'^\*\s+\*\*Visual Suggestion\*\*:\s*(.*)', line_stripped, re.IGNORECASE)
            if visual_match:
                val = visual_match.group(1).strip()
                if val:
                    slide_info['visual_suggestions'].append(val)
                current_section = 'visual'
                continue

            script_match = re.match(r'^\*\s+\*\*Presenter Script\*\*:\s*(.*)', line_stripped, re.IGNORECASE)
            if script_match:
                val = script_match.group(1).strip()
                if val:
                    val = re.sub(r'^>\s*', '', val)
                    slide_info['presenter_script'].append(val)
                current_section = 'script'
                continue

            checkpoint_match = re.match(r'^\*\s+\*\*Interactive Checkpoint\*\*:\s*(.*)', line_stripped, re.IGNORECASE)
            if checkpoint_match:
                val = checkpoint_match.group(1).strip()
                if val:
                    slide_info['interactive_checkpoints'].append(val)
                current_section = 'checkpoint'
                continue

            # If it starts with another bullet * **Name**:
            other_bullet_match = re.match(r'^\*\s+\*\*(.*?)\*\*:\s*(.*)', line_stripped)
            if other_bullet_match:
                val = other_bullet_match.group(2).strip()
                if val:
                    slide_info['other_bullets'].append(f"{other_bullet_match.group(1)}: {val}")
                else:
                    slide_info['other_bullets'].append(other_bullet_match.group(1))
                current_section = 'other'
                continue

            # Appending items to current section
            if current_section == 'visual':
                cleaned = re.sub(r'^[\*\-\+]\s*', '', line_stripped)
                slide_info['visual_suggestions'].append(cleaned)
            elif current_section == 'script':
                cleaned = re.sub(r'^>\s*', '', line_stripped)
                cleaned = re.sub(r'^[\*\-\+]\s*', '', cleaned)
                slide_info['presenter_script'].append(cleaned)
            elif current_section == 'checkpoint':
                cleaned = re.sub(r'^[\*\-\+]\s*', '', line_stripped)
                slide_info['interactive_checkpoints'].append(cleaned)
            elif current_section == 'other':
                cleaned = re.sub(r'^[\*\-\+]\s*', '', line_stripped)
                slide_info['other_bullets'].append(cleaned)
            else:
                cleaned = re.sub(r'^[\*\-\+]\s*', '', line_stripped)
                slide_info['other_bullets'].append(cleaned)

        slide_info['title'] = slide_info['title'] or slide_info['header']
        slides_data.append(slide_info)

    return slides_data

def prepare_assets(assets_dir):
    os.makedirs(assets_dir, exist_ok=True)
    mapped_images = {}
    
    for slide_idx, filename in IMAGE_MAPPING.items():
        local_path = os.path.join(assets_dir, f"slide_{slide_idx}.png")
        artifact_path = os.path.join(ARTIFACT_DIR, filename)
        
        # Copy from artifact directory to local asset folder if it exists
        if os.path.exists(artifact_path):
            shutil.copy(artifact_path, local_path)
            mapped_images[slide_idx] = local_path
        elif os.path.exists(local_path):
            mapped_images[slide_idx] = local_path
        else:
            print(f"Warning: Missing image for slide {slide_idx} ({filename})")
            
    return mapped_images

def apply_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def create_presentation(slides_data, images_map, output_path):
    prs = Presentation()
    
    # 16:9 Widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Use slide layout 6 (Blank) for full styling control
    blank_layout = prs.slide_layouts[6]

    for idx, slide_info in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        apply_background(slide)
        
        slide_num = idx + 1
        
        # --- Left Column: Text & Content ---
        if slide_num == 1:
            # Welcome slide
            # Decorative Purple Line
            line = slide.shapes.add_shape(
                1, # MSO_SHAPE.RECTANGLE
                Inches(0.8), Inches(1.8), Inches(0.12), Inches(3.6)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = ACCENT_PURPLE
            line.line.color.rgb = ACCENT_PURPLE
            
            # Text box
            txBox = slide.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(6.0), Inches(3.6))
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0)
            
            # Title
            p_title = tf.paragraphs[0]
            p_title.text = slide_info['title']
            p_title.font.name = "Segoe UI"
            p_title.font.size = Pt(44)
            p_title.font.bold = True
            p_title.font.color.rgb = TEXT_WHITE
            p_title.space_after = Pt(20)
            
            # Subtitle
            if slide_info['subtitle']:
                p_sub = tf.add_paragraph()
                p_sub.text = slide_info['subtitle']
                p_sub.font.name = "Segoe UI"
                p_sub.font.size = Pt(20)
                p_sub.font.color.rgb = ACCENT_BLUE
                
        else:
            # Normal slide
            # Header
            txHeader = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8))
            tfHeader = txHeader.text_frame
            tfHeader.word_wrap = True
            tfHeader.margin_left = tfHeader.margin_top = tfHeader.margin_right = tfHeader.margin_bottom = Inches(0)
            
            p_head = tfHeader.paragraphs[0]
            p_head.text = slide_info['title']
            p_head.font.name = "Segoe UI"
            p_head.font.size = Pt(28)
            p_head.font.bold = True
            p_head.font.color.rgb = TEXT_WHITE
            
            # Horizontal line separator
            line = slide.shapes.add_shape(
                1, # RECTANGLE
                Inches(0.8), Inches(1.2), Inches(11.733), Inches(0.02)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = ACCENT_BLUE
            line.line.color.rgb = ACCENT_BLUE
            
            # Main Left Body content
            txBody = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(6.0), Inches(5.0))
            tf = txBody.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = Inches(0)
            
            first = True
            if slide_info['interactive_checkpoints']:
                p = tf.paragraphs[0]
                p.text = "Interactive Activity:"
                p.font.name = "Segoe UI"
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = ACCENT_PURPLE
                p.space_after = Pt(8)
                first = False
                
                for item in slide_info['interactive_checkpoints']:
                    p = tf.add_paragraph()
                    p.text = f"• {item}"
                    p.font.name = "Segoe UI"
                    p.font.size = Pt(15)
                    p.font.color.rgb = TEXT_LIGHT_GRAY
                    p.space_after = Pt(6)
            
            if slide_info['other_bullets']:
                if not first:
                    p = tf.add_paragraph()
                    p.space_after = Pt(12)
                else:
                    p = tf.paragraphs[0]
                    first = False
                
                for item in slide_info['other_bullets']:
                    # Highlight bold texts inside bullet points if any
                    p = tf.add_paragraph()
                    p.text = f"• {item}"
                    p.font.name = "Segoe UI"
                    p.font.size = Pt(15)
                    p.font.color.rgb = TEXT_LIGHT_GRAY
                    p.space_after = Pt(6)
                    
            if first:
                # Placeholder content if no bullet text exists
                p = tf.paragraphs[0]
                p.text = "Let's explore this topic together!"
                p.font.name = "Segoe UI"
                p.font.size = Pt(18)
                p.font.italic = True
                p.font.color.rgb = TEXT_LIGHT_GRAY

        # --- Right Column: Image block ---
        if slide_num in images_map:
            img_path = images_map[slide_num]
            # Coordinates: Right half, centered vertically
            left = Inches(7.5)
            top = Inches(1.6)
            width = Inches(5.0)
            height = Inches(4.5)
            slide.shapes.add_picture(img_path, left, top, width, height)

        # --- Slide Notes (Visual Suggestions & Presenter Script) ---
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        
        notes_content = []
        if slide_info['visual_suggestions']:
            notes_content.append("--- VISUAL SUGGESTIONS ---")
            notes_content.extend([f"- {s}" for s in slide_info['visual_suggestions']])
            notes_content.append("")
            
        if slide_info['presenter_script']:
            notes_content.append("--- PRESENTER SCRIPT ---")
            notes_content.extend(slide_info['presenter_script'])
            notes_content.append("")
            
        if slide_info['interactive_checkpoints']:
            notes_content.append("--- INTERACTIVE CHECKPOINTS ---")
            notes_content.extend([f"- {c}" for c in slide_info['interactive_checkpoints']])
            notes_content.append("")

        notes_text_frame.text = "\n".join(notes_content).strip()

    try:
        prs.save(output_path)
        print(f"Successfully generated PowerPoint presentation at: {output_path}")
    except PermissionError:
        print("\n" + "="*60)
        print(f"ERROR: Permission denied when saving to '{output_path}'.")
        print("This usually means the file is open in Microsoft PowerPoint or another application.")
        print("Please CLOSE the PowerPoint window and run this script again!")
        print("="*60 + "\n")
        sys.exit(1)

if __name__ == "__main__":
    input_file = os.path.join("module_01_intro_to_ai", "slides.md")
    output_file = os.path.join("module_01_intro_to_ai", "slides.pptx")
    assets_folder = os.path.join("module_01_intro_to_ai", "assets")
    
    if not os.path.exists(input_file):
        print(f"Error: Could not find input file: {input_file}")
        sys.exit(1)
        
    slides_data = parse_markdown_slides(input_file)
    images_map = prepare_assets(assets_folder)
    create_presentation(slides_data, images_map, output_file)
